#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Image-centric v6 comparison presentation deck (fallback / portable artifact).

Builds:
  1. results/vis/_slides_assets_v6/metrics_bar.png   (matplotlib SSIM/PSNR bars)
  2. 발표 자료/ViT_v6_비교_이미지중심.pptx        (python-pptx, 16:9)

Run with the conda `mri_env` python (has python-pptx, PIL, matplotlib):
    python tools/build_v6_deck.py

This is the FALLBACK to the Figma Slides plan (see docs/figma_deck_handoff.md).
The deck is image-centric and built around results/vis/root_track/compare_versions_aligned/.
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSET = os.path.join(ROOT, "results/vis/_slides_assets_v6")
os.makedirs(ASSET, exist_ok=True)
CHART = os.path.join(ASSET, "metrics_bar.png")
OUT = os.path.join(ROOT, "발표 자료", "ViT_v6_비교_이미지중심.pptx")

# ---- full-eval metrics (7,270 slices, R=4, skimage standard SSIM) -----------
# source: docs/presentation_overview.md §6.1
MODELS = ["U-Net\n(496M)", "SS2D v6", "ETER v6", "SS2D v6_3"]
SSIM = [0.8858, 0.8913, 0.8862, 0.8924]
PSNR = [34.66, 35.96, 34.63, 36.05]
COLORS = ["#9AA3AD", "#2E6FB5", "#2E8B57", "#16406E"]  # gray, blue, green, dark-blue

# ---------------------------------------------------------------------------
# 1) metrics bar chart
# ---------------------------------------------------------------------------
def build_chart():
    fig, axes = plt.subplots(1, 2, figsize=(11, 4.2), dpi=160)
    for ax, vals, title, lo, hi, fmt in [
        (axes[0], SSIM, "SSIM (higher = better)", 0.875, 0.895, "{:.4f}"),
        (axes[1], PSNR, "PSNR / dB (higher = better)", 33.5, 36.6, "{:.2f}"),
    ]:
        bars = ax.bar(MODELS, vals, color=COLORS, width=0.62, edgecolor="white")
        ax.set_ylim(lo, hi)
        ax.set_title(title, fontsize=13, fontweight="bold", color="#1F2D3D")
        ax.spines[["top", "right"]].set_visible(False)
        ax.tick_params(labelsize=10)
        ax.grid(axis="y", color="#E3E7EC", linewidth=0.8)
        ax.set_axisbelow(True)
        for b, v in zip(bars, vals):
            ax.text(b.get_x() + b.get_width() / 2, v, fmt.format(v),
                    ha="center", va="bottom", fontsize=10, fontweight="bold",
                    color="#1F2D3D")
    # U-Net baseline reference line on SSIM
    axes[0].axhline(SSIM[0], color="#9AA3AD", ls="--", lw=1)
    fig.suptitle("Full evaluation (7,270 slices · R=4)  —  v6 overtakes U-Net",
                 fontsize=14, fontweight="bold", color="#16406E", y=1.02)
    fig.tight_layout()
    fig.savefig(CHART, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print("chart ->", CHART)

# ---------------------------------------------------------------------------
# 2) pptx deck
# ---------------------------------------------------------------------------
def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    NAVY = RGBColor(0x16, 0x40, 0x6E)
    DARK = RGBColor(0x1F, 0x2D, 0x3D)
    BLUE = RGBColor(0x2E, 0x6F, 0xB5)
    GREEN = RGBColor(0x2E, 0x8B, 0x57)
    GRAY = RGBColor(0x5C, 0x66, 0x70)
    LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    KFONT = "Apple SD Gothic Neo"  # macOS Korean font; Keynote/PPT fall back gracefully

    def slide(bg=WHITE):
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg
        return s

    def box(s, l, t, w, h):
        tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        return tf

    def setp(p, text, size, color, bold=False, align=PP_ALIGN.LEFT, font=KFONT):
        p.text = text
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
        return p

    def band(s, t, h, color):
        from pptx.enum.shapes import MSO_SHAPE
        shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(t), SW, Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def header(s, title, sub=None):
        band(s, 0, 1.0, NAVY)
        tf = box(s, 0.55, 0.16, 12.2, 0.7)
        setp(tf.paragraphs[0], title, 26, WHITE, bold=True)
        if sub:
            tf2 = box(s, 0.57, 0.78, 12.2, 0.35)
            setp(tf2.paragraphs[0], sub, 13, RGBColor(0xC8, 0xD6, 0xE8))

    def footer(s):
        tf = box(s, 0.55, 7.05, 12.2, 0.35)
        setp(tf.paragraphs[0],
             "fastMRI brain AXFLAIR multicoil · R=4 equispaced · 8GB GPU · 2026-06-01",
             9, GRAY)

    def picture_fit(s, path, l, t, maxw, maxh):
        """place image preserving aspect ratio inside box (l,t,maxw,maxh) in inches, centered."""
        from PIL import Image
        iw, ih = Image.open(path).size
        ar = iw / ih
        boxar = maxw / maxh
        if ar > boxar:
            w = maxw; h = maxw / ar
        else:
            h = maxh; w = maxh * ar
        ox = l + (maxw - w) / 2
        oy = t + (maxh - h) / 2
        return s.shapes.add_picture(path, Inches(ox), Inches(oy), Inches(w), Inches(h))

    # -- 1. TITLE -----------------------------------------------------------
    s = slide(NAVY)
    tf = box(s, 0.9, 2.5, 11.5, 2.0)
    setp(tf.paragraphs[0], "ViT 기반 fastMRI 재구성", 40, WHITE, bold=True)
    p = tf.add_paragraph(); setp(p, "v6 결과 비교 — SS2D(Mamba) vs ETER(GRU) vs U-Net", 24,
                                 RGBColor(0x9D, 0xC2, 0xEE))
    tf2 = box(s, 0.92, 4.6, 11.5, 1.0)
    setp(tf2.paragraphs[0], "이미지 중심 비교본  ·  정렬 4×4 그리드(재구성 + 오차맵)", 15,
         RGBColor(0xC8, 0xD6, 0xE8))
    p = tf2.add_paragraph(); setp(p, "표준 skimage SSIM · 전체 7,270 슬라이스 평가", 15,
                                  RGBColor(0xC8, 0xD6, 0xE8))

    # -- 2. HOW TO READ -----------------------------------------------------
    s = slide()
    header(s, "그림 읽는 법 — 4×4 비교 그리드", "각 PNG 한 장 = 한 슬라이스의 전체 모델 비교")
    picture_fit(s, os.path.join(ASSET, "v6cmp_1982.png"), 0.4, 1.25, 6.1, 5.6)
    tf = box(s, 6.9, 1.45, 6.0, 5.4)
    rows = [
        ("1행 (재구성)", "GT │ SS2D v4 │ SS2D v6 │ SS2D v6_1", BLUE),
        ("2행 (오차맵)", "위 SS2D 모델들의 |예측 − GT|", BLUE),
        ("3행 (재구성)", "U-Net(pretrained) │ ETER v4 │ ETER v6 │ ETER v6_1", GREEN),
        ("4행 (오차맵)", "위 ETER 모델들의 |예측 − GT|", GREEN),
    ]
    first = True
    for head, body, col in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        setp(p, head, 16, col, bold=True); p.space_after = Pt(2)
        p2 = tf.add_paragraph(); setp(p2, "   " + body, 14, DARK); p2.space_after = Pt(10)
    for line in [
        "· 패널 상단 숫자 = 해당 슬라이스의 PSNR / SSIM (raw amplitude)",
        "· 오차맵이 어두울수록(빨강이 약할수록) 오차가 작음",
        "· 동일 색상 척도(err_vmax)로 정렬되어 모델 간 직접 비교 가능",
    ]:
        p = tf.add_paragraph(); setp(p, line, 13, GRAY); p.space_after = Pt(4)
    footer(s)

    # -- 3..N. IMAGE SLIDES -------------------------------------------------
    img_slides = [
        ("v6cmp_0660.png", "슬라이스 #0660 — 중간 뇌 (뇌실 레벨)", [
            "v4 → v6에서 뇌실·실질 경계가 또렷해짐.",
            "SS2D 오차맵이 ETER보다 전반적으로 옅음 → SS2D 우위.",
            "v6_1(gradient loss)은 미세 과샤프닝 경향.",
        ]),
        ("v6cmp_1982.png", "슬라이스 #1982 — 중간 뇌", [
            "구조 복원의 일관성을 보여주는 대표 슬라이스.",
            "SS2D v6의 오차맵이 ETER v6보다 균일하게 옅음.",
            "U-Net 대비 v6의 PSNR/SSIM 동등~우위.",
        ]),
        ("v6cmp_3304.png", "슬라이스 #3304 — 두정부 (sulci 미세구조)", [
            "SS2D v6 SSIM 0.9367 > ETER v6 0.9242 (U-Net 0.9447).",
            "v6_1: PSNR 33.71 → 32.48 dB 하락 = over-sharpening 부작용.",
            "고주파 sulci 영역에서 모델 차이가 가장 잘 드러남.",
        ]),
        ("v6cmp_5286.png", "슬라이스 #5286 — 두정부 고주파 (어려운 케이스)", [
            "sulci·혈관 등 fine detail에 오차가 집중.",
            "L1+SSIM loss의 mean-prediction blurring이 가시화 (발견 ②).",
            "정량 지표는 높아도 시각적으로 흐릿한 괴리의 근거.",
        ]),
        ("v6cmp_6608.png", "슬라이스 #6608 — 두정부 고detail", [
            "여러 해부 레벨에서 SS2D > ETER가 일관되게 관찰됨.",
            "4-방향 SSM의 장거리 의존성 포착이 k-space aliasing에 효과적.",
        ]),
    ]
    for fname, title, bullets in img_slides:
        s = slide()
        header(s, title)
        path = os.path.join(ASSET, fname)
        if os.path.exists(path):
            picture_fit(s, path, 0.35, 1.2, 8.0, 5.7)
        tf = box(s, 8.65, 1.55, 4.35, 5.2)
        first = True
        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            setp(p, "• " + b, 15, DARK); p.space_after = Pt(12)
        footer(s)

    # -- SUMMARY ------------------------------------------------------------
    s = slide()
    header(s, "정량 요약 — 표준 SSIM 전체 평가 (7,270 슬라이스)")
    if os.path.exists(CHART):
        picture_fit(s, CHART, 0.4, 1.25, 7.4, 4.3)
    # table
    from pptx.util import Inches as IN
    rows_t = [
        ("모델", "PSNR(dB)", "SSIM", "NMSE"),
        ("U-Net (pretrained, 496M)", "34.66", "0.8858", "0.00737"),
        ("SS2D-ViT v6  ▲U-Net 추월", "35.96", "0.8913", "0.00810"),
        ("ETER-ViT v6", "34.63", "0.8862", "0.01080"),
        ("SS2D-ViT v6_3  ★채택후보", "36.05", "0.8924", "0.00800"),
    ]
    tbl = s.shapes.add_table(len(rows_t), 4, IN(0.5), IN(5.7), IN(12.3), IN(1.3)).table
    tbl.columns[0].width = IN(6.0)
    for c in (1, 2, 3):
        tbl.columns[c].width = IN(2.1)
    for ri, row in enumerate(rows_t):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            for r in para.runs:
                r.font.size = Pt(12)
                r.font.name = KFONT
                r.font.bold = (ri == 0)
                r.font.color.rgb = WHITE if ri == 0 else DARK
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif "v6_3" in row[0]:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xDD, 0xEA, 0xF7)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
    # right-side note
    tf = box(s, 8.05, 1.4, 4.9, 4.0)
    setp(tf.paragraphs[0], "핵심", 16, NAVY, bold=True)
    for line in [
        "표준 metric에서 v6가 U-Net(0.8858)을",
        "SSIM·PSNR 모두에서 추월.",
        "",
        "SS2D(Mamba) > ETER(GRU) 일관.",
        "",
        "v6_3(정규화 완화)가 모든 지표에서",
        "v6보다 개선 → 채택 후보.",
    ]:
        p = tf.add_paragraph(); setp(p, line, 13, DARK)
    footer(s)

    # -- CONCLUSION ---------------------------------------------------------
    s = slide()
    header(s, "핵심 메시지")
    tf = box(s, 0.7, 1.5, 11.9, 5.2)
    pts = [
        ("1. 표준 SSIM 기준 v6가 U-Net 추월",
         "SS2D v6 SSIM 0.8913 / PSNR 35.96 dB > U-Net 0.8858 / 34.66 dB."),
        ("2. Mamba(SS2D) > Bi-GRU(ETER) 일관",
         "동일 dataloader·레시피에서 SSIM·PSNR·오차맵 모두 SS2D 우위."),
        ("3. raw SSIM 배경 부풀림 + mean-prediction blurring (발견 ②)",
         "슬라이스 50%+가 배경이라 raw SSIM이 부풀려짐. 시각적 흐림의 본질은 L1+SSIM loss."),
        ("4. v6_3 채택 후보",
         "정규화 완화(dropout 0.2→0.1, WD 3e-5→1e-5)로 모든 지표 개선: SSIM 0.8924 / PSNR 36.05 dB."),
    ]
    first = True
    for head, body in pts:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        setp(p, head, 19, NAVY, bold=True); p.space_before = Pt(8); p.space_after = Pt(2)
        p2 = tf.add_paragraph(); setp(p2, "    " + body, 14, DARK); p2.space_after = Pt(8)
    footer(s)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("pptx ->", OUT, "(slides:", len(prs.slides._sldIdLst), ")")


if __name__ == "__main__":
    build_chart()
    build_pptx()
